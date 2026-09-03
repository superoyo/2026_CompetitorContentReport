# -*- coding: utf-8 -*-
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from report_config import BRANDS, ANALYSIS

# ---------- palette ----------
DARK = RGBColor(0x0B, 0x14, 0x22)
LIGHT_TINT = RGBColor(0xEE, 0xF2, 0xF9)
ROW_ALT = RGBColor(0xF6, 0xF8, 0xFC)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xFC, 0xA3, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x16, 0x1F, 0x30)
MUTED = RGBColor(0x64, 0x74, 0x88)
RING = RGBColor(0x1E, 0x30, 0x50)
PANEL = RGBColor(0x11, 0x1A, 0x2E)

HEAD_FONT = "Prompt"
BODY_FONT = "Sarabun"

def hx(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def tint(h, factor=0.90):
    """lighten a hex color toward white by factor (0..1)."""
    h = h.lstrip('#')
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(r + (255 - r) * factor); g = int(g + (255 - g) * factor); b = int(b + (255 - b) * factor)
    return RGBColor(r, g, b)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

P = json.load(open('/tmp/processed_8.json'))
AGG = P['agg']; MET = P['metrics']; TOP5 = P['top5']

NAME = {b[0]: b[1] for b in BRANDS}
LETTER = {b[0]: b[2] for b in BRANDS}
COLOR = {b[0]: b[3] for b in BRANDS}
ORDER = sorted([b[0] for b in BRANDS], key=lambda k: AGG[k]['total'], reverse=True)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(blank)

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def no_line(shape):
    shape.line.fill.background()

def add_shadow(shp, blur=70000, dist=24000, alpha=20000):
    el = shp._element.spPr
    for ex in el.findall(qn('a:effectLst')):
        el.remove(ex)
    eff = el.makeelement(qn('a:effectLst'), {})
    sh = el.makeelement(qn('a:outerShdw'), {'blurRad': str(blur), 'dist': str(dist), 'dir': '5400000', 'rotWithShape': '0'})
    clr = el.makeelement(qn('a:srgbClr'), {'val': '000000'})
    a = el.makeelement(qn('a:alpha'), {'val': str(alpha)})
    clr.append(a); sh.append(clr); eff.append(sh); el.append(eff)

def add_rect(slide, x, y, w, h, fill=None, shadow=False, radius=None, line_color=None, line_w=1.0):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line_color is None:
        no_line(shp)
    else:
        shp.line.color.rgb = line_color; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow: add_shadow(shp)
    return shp

def add_ring(slide, cx, cy, d, color, w=1.4):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - d/2), Inches(cy - d/2), Inches(d), Inches(d))
    shp.fill.background(); shp.line.color.rgb = color; shp.line.width = Pt(w); shp.shadow.inherit = False
    return shp

def add_badge(slide, x, y, d, fill, text, tc=WHITE, fs=14):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill; no_line(shp); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = False
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf, m, 0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.size = Pt(fs); r.font.bold = True
    r.font.color.rgb = tc; r.font.name = HEAD_FONT
    return shp

def add_text(slide, x, y, w, h, text, size=14, color=TEXT_DARK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=BODY_FONT, ls=None, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf, m, 0)
    tf.vertical_anchor = anchor
    if isinstance(text, str): text = [(text, {})]
    first = True
    for t, o in text:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        p.alignment = o.get('align', align)
        if ls: p.line_spacing = ls
        if o.get('space_before'): p.space_before = Pt(o['space_before'])
        r = p.add_run(); r.text = t; f = r.font
        f.size = Pt(o.get('size', size)); f.bold = o.get('bold', bold); f.italic = o.get('italic', italic)
        f.name = o.get('font', font); f.color.rgb = o.get('color', color)
    return tb

def truncate(s, n):
    s = (s or "").strip().replace('\n', ' ')
    return s if len(s) <= n else s[:n].rstrip() + "…"

def fmt(n): return f"{n:,}"

def img_for(key, i):
    p = f"post_images_cropped/{key}_{i}.jpg"
    return p if os.path.exists(p) else None

# ============ Slide 1: Title ============
s = add_slide(); set_bg(s, DARK)
add_ring(s, 13.7, -0.3, 2.6, RING, 1.6); add_ring(s, 13.0, 0.7, 1.1, RING, 1.6); add_ring(s, -0.3, 8.0, 1.8, RING, 1.6)
add_rect(s, 1.0, 2.05, 0.9, 0.09, fill=ACCENT)
add_text(s, 1.0, 2.25, 11.3, 0.5, "รายงานสรุป Engagement บน Facebook", size=19, color=hx('#9FC9DE'), bold=True, font=HEAD_FONT)
add_text(s, 1.0, 2.72, 11.5, 1.3, "Top 5 คอนเทนต์ยอด Engagement สูงสุด", size=42, color=WHITE, bold=True, font=HEAD_FONT)
add_text(s, 1.0, 3.7, 11.3, 0.6, "ประจำเดือนพฤษภาคม 2569  (May 2026)  ·  8 เพจ", size=19, color=hx('#CFE7F0'), font=HEAD_FONT)

bx, by = 1.0, 4.75
for key in [b[0] for b in BRANDS]:
    add_badge(s, bx, by, 0.5, hx(COLOR[key]), LETTER[key], fs=15)
    add_text(s, bx - 0.62, by + 0.56, 1.74, 0.5, NAME[key].replace(' Thailand', ''), size=10.5,
             color=hx('#CFE7F0'), align=PP_ALIGN.CENTER, wrap=True, font=BODY_FONT)
    bx += 1.5
add_text(s, 1.0, 6.9, 9.0, 0.4, "ข้อมูลจากโพสต์สาธารณะบนเพจ Facebook · ดึงผ่าน Apify", size=11, color=hx('#7C9CB0'))

# ============ Slide 2: Overview ============
s = add_slide(); set_bg(s, WHITE)
add_text(s, 0.6, 0.42, 10, 0.6, "ภาพรวม Engagement เดือนพฤษภาคม 2569", size=27, color=hx('#0B2545'), bold=True, font=HEAD_FONT)
add_text(s, 0.6, 1.0, 10, 0.4, "เปรียบเทียบยอด Engagement รวม (Likes + Comments + Shares) ทั้ง 8 เพจ", size=13, color=MUTED)

grand = sum(AGG[k]['total'] for k in AGG)
posts = sum(AGG[k]['posts'] for k in AGG)
add_rect(s, 10.5, 0.42, 2.3, 1.12, fill=LIGHT_TINT, radius=0.16, shadow=True)
add_text(s, 10.5, 0.55, 2.3, 0.5, fmt(grand), size=25, color=hx('#0B2545'), bold=True, align=PP_ALIGN.CENTER, font=HEAD_FONT)
add_text(s, 10.5, 1.06, 2.3, 0.4, f"Engagement รวม · {posts} โพสต์", size=10, color=MUTED, align=PP_ALIGN.CENTER)

cd = CategoryChartData()
cd.categories = [NAME[k].replace(' Thailand', '') for k in ORDER]
cd.add_series("Total", [AGG[k]['total'] for k in ORDER])
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(1.75), Inches(7.2), Inches(4.9), cd)
ch = gf.chart; ch.has_legend = False; ch.has_title = False
plot = ch.plots[0]; plot.has_data_labels = True
dl = plot.data_labels; dl.number_format = '#,##0'; dl.number_format_is_linked = False
dl.font.size = Pt(11); dl.font.bold = True; dl.font.color.rgb = TEXT_DARK; dl.font.name = BODY_FONT; dl.position = 2
ser = plot.series[0]
ser.format.fill.solid(); ser.format.fill.fore_color.rgb = hx('#94A3B8')
for i, pt in enumerate(ser.points):
    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = hx(COLOR[ORDER[i]])
ca = ch.category_axis; ca.tick_labels.font.size = Pt(10.5); ca.tick_labels.font.color.rgb = TEXT_DARK
ca.tick_labels.font.name = BODY_FONT; ca.format.line.color.rgb = hx('#D8E3E8')
va = ch.value_axis; va.visible = False; va.has_major_gridlines = False

# ranking rows
rx, rw, rh, rgap, ry = 8.25, 4.55, 0.60, 0.105, 1.75
for rank, key in enumerate(ORDER, 1):
    v = AGG[key]
    add_rect(s, rx, ry, rw, rh, fill=(tint(COLOR[key], 0.86) if rank == 1 else ROW_ALT), radius=0.16)
    add_badge(s, rx + 0.12, ry + (rh - 0.44) / 2, 0.44, hx(COLOR[key]), str(rank), fs=13)
    add_text(s, rx + 0.68, ry + 0.06, 2.2, 0.3, NAME[key].replace(' Thailand', ''), size=12.5, color=TEXT_DARK, bold=True, font=HEAD_FONT)
    add_text(s, rx + 0.68, ry + 0.33, 2.2, 0.24, f"{v['posts']} โพสต์ · เฉลี่ย {fmt(round(v['avg']))}/โพสต์", size=8.5, color=MUTED)
    add_text(s, rx + 2.7, ry, 1.75, rh, fmt(v['total']), size=16, color=hx(COLOR[key]), bold=True,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, font=HEAD_FONT)
    ry += rh + rgap
add_text(s, 0.6, 6.95, 11.5, 0.35, "Engagement = Likes/Reactions + Comments + Shares (ข้อมูลจากโพสต์สาธารณะ)", size=10, color=MUTED, italic=True)

# ============ per-brand Top5 + Analysis ============
CARD_W, CARD_GAP, IMG_H, X0, CARD_TOP = 2.15, 0.35, 2.62, 0.6, 1.62

for key in ORDER:
    display = NAME[key]; c = hx(COLOR[key]); v = AGG[key]
    # --- Top5 slide ---
    s = add_slide(); set_bg(s, WHITE)
    add_badge(s, 0.6, 0.42, 0.55, c, LETTER[key], fs=20)
    add_text(s, 1.32, 0.4, 8.6, 0.5, display, size=24, color=hx('#0B2545'), bold=True, font=HEAD_FONT)
    add_text(s, 1.32, 0.92, 8.6, 0.4, "Top 5 คอนเทนต์ Engagement สูงสุด — พฤษภาคม 2569", size=12.5, color=MUTED)
    add_rect(s, 10.45, 0.35, 2.3, 1.05, fill=tint(COLOR[key], 0.86), radius=0.16, shadow=True)
    add_text(s, 10.45, 0.46, 2.3, 0.46, fmt(v['total']), size=22, color=c, bold=True, align=PP_ALIGN.CENTER, font=HEAD_FONT)
    add_text(s, 10.45, 0.94, 2.3, 0.34, f"Total Engagement · {v['posts']} โพสต์", size=9, color=MUTED, align=PP_ALIGN.CENTER)

    cx = X0
    for i, post in enumerate(TOP5[key][:5], 1):
        card_h = 4.95
        add_rect(s, cx, CARD_TOP, CARD_W, card_h, fill=CARD_BG, radius=0.05, shadow=True)
        ix, iy, iw, ih = cx + 0.06, CARD_TOP + 0.06, CARD_W - 0.12, IMG_H - 0.06
        ip = img_for(key, i)
        if ip:
            s.shapes.add_picture(ip, Inches(ix), Inches(iy), Inches(iw), Inches(ih))
        else:
            add_rect(s, ix, iy, iw, ih, fill=tint(COLOR[key], 0.85), radius=0.04)
            add_text(s, ix, iy, iw, ih, [("“", {"size": 30, "color": c, "bold": True, "align": PP_ALIGN.CENTER}),
                                          ("โพสต์ข้อความ", {"size": 11, "color": MUTED, "bold": True, "align": PP_ALIGN.CENTER})],
                     anchor=MSO_ANCHOR.MIDDLE)
        add_badge(s, cx - 0.03, CARD_TOP - 0.03, 0.44, (ACCENT if i == 1 else c), str(i), fs=15)
        ty = CARD_TOP + IMG_H + 0.06
        add_text(s, cx + 0.1, ty, CARD_W - 0.2, 0.22, (post['time'] or '')[:10], size=8.5, color=MUTED)
        add_text(s, cx + 0.1, ty + 0.22, CARD_W - 0.2, 0.72, truncate(post['text'], 72), size=8.5, color=TEXT_DARK, ls=1.05)
        sy = ty + 1.02; sw = (CARD_W - 0.2) / 3
        for j, (lab, val) in enumerate([("Likes", post['likes']), ("Comments", post['comments']), ("Shares", post['shares'])]):
            sx = cx + 0.1 + j * sw
            add_text(s, sx, sy, sw, 0.28, fmt(val), size=12, color=c, bold=True, align=PP_ALIGN.CENTER, font=HEAD_FONT)
            add_text(s, sx, sy + 0.27, sw, 0.22, lab, size=7, color=MUTED, align=PP_ALIGN.CENTER)
        byy = sy + 0.56
        add_rect(s, cx + 0.1, byy, CARD_W - 0.2, 0.42, fill=(ACCENT if i == 1 else tint(COLOR[key], 0.80)), radius=0.22)
        add_text(s, cx + 0.1, byy, CARD_W - 0.2, 0.42, fmt(post['total']) + " Total",
                 size=11.5, color=(DARK if i == 1 else c), bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD_FONT)
        cx += CARD_W + CARD_GAP
    add_text(s, 0.6, CARD_TOP + 4.95 + 0.12, 12.1, 0.3, "Total = Likes/Reactions + Comments + Shares", size=9.5, color=MUTED, italic=True)

    # --- Analysis + Reco slide ---
    s = add_slide(); set_bg(s, DARK)
    add_badge(s, 0.6, 0.5, 0.55, c, LETTER[key], fs=20)
    add_text(s, 1.32, 0.46, 9.5, 0.5, display, size=23, color=WHITE, bold=True, font=HEAD_FONT)
    add_text(s, 1.32, 0.98, 10.5, 0.4, f"บทวิเคราะห์คอนเทนต์ & ข้อเสนอแนะเดือนถัดไป — วิเคราะห์จาก {v['posts']} โพสต์ในเดือน พ.ค.",
             size=12, color=hx('#9FC9DE'))
    # chips
    a = ANALYSIS[key]; chx = 0.6
    for chip in a['chips']:
        w_est = 0.16 + len(chip) * 0.088
        add_rect(s, chx, 1.6, w_est, 0.42, fill=PANEL, radius=0.3, line_color=hx('#2A3A57'), line_w=1.0)
        add_text(s, chx, 1.6, w_est, 0.42, chip, size=10.5, color=hx('#C7D6EA'), bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        chx += w_est + 0.18
    # two boxes
    box_y, box_w, box_h = 2.35, 6.0, 4.5
    # analysis box
    add_rect(s, 0.6, box_y, box_w, box_h, fill=PANEL, radius=0.05, line_color=hx('#26507A'), line_w=1.25)
    add_rect(s, 0.6, box_y, 0.12, box_h, fill=hx('#38BDF8'))
    add_text(s, 0.95, box_y + 0.25, box_w - 0.6, 0.4, "📊  บทวิเคราะห์คอนเทนต์", size=15, color=hx('#7DD3FC'), bold=True, font=HEAD_FONT)
    ay = box_y + 0.9
    for pt in a['analysis']:
        add_text(s, 0.95, ay, 0.3, 0.4, "▸", size=13, color=hx('#38BDF8'), bold=True)
        tb = add_text(s, 1.3, ay, box_w - 1.0, 1.2, pt, size=11.5, color=hx('#DCE6F2'), ls=1.12)
        ay += 1.16
    # reco box
    rxb = 0.6 + box_w + 0.35
    add_rect(s, rxb, box_y, box_w, box_h, fill=PANEL, radius=0.05, line_color=hx('#7A5A1E'), line_w=1.25)
    add_rect(s, rxb, box_y, 0.12, box_h, fill=ACCENT)
    add_text(s, rxb + 0.35, box_y + 0.25, box_w - 0.6, 0.4, "🚀  ควรทำต่อในเดือนถัดไป", size=15, color=hx('#FCD34D'), bold=True, font=HEAD_FONT)
    ry2 = box_y + 0.9
    for pt in a['reco']:
        add_text(s, rxb + 0.35, ry2, 0.3, 0.4, "✓", size=13, color=hx('#FBBF24'), bold=True)
        add_text(s, rxb + 0.72, ry2, box_w - 1.0, 1.2, pt, size=11.5, color=hx('#F1E7D2'), ls=1.12)
        ry2 += 1.16

# ============ Closing ============
s = add_slide(); set_bg(s, DARK)
add_ring(s, 13.6, 7.9, 2.6, RING, 1.6); add_ring(s, -0.3, -0.3, 1.3, RING, 1.6)
add_text(s, 1.0, 1.0, 11.3, 0.6, "หมายเหตุข้อมูล", size=26, color=WHITE, bold=True, font=HEAD_FONT)
add_rect(s, 1.0, 1.72, 0.75, 0.06, fill=ACCENT)
notes = [
    "ข้อมูลดึงจากโพสต์สาธารณะบนเพจ Facebook ผ่านเครื่องมือสแครปข้อมูล (Apify) ไม่ใช่ตัวเลขจาก Facebook Page Insights โดยตรง",
    "Engagement นับจาก Likes/Reactions + Comments + Shares ของแต่ละโพสต์ ไม่รวม Reach, Impressions หรือ Click ซึ่งดูได้จาก Insights เท่านั้น",
    "ช่วงข้อมูล: 1–31 พฤษภาคม 2569 (May 2026) · จำนวนโพสต์ที่ดึงได้ต่อเพจอาจต่างกันตามความถี่การโพสต์จริง",
    "เพจที่มีจำนวนโพสต์น้อย ตัวเลขจึงสะท้อนช่วงตัวอย่างที่จำกัด",
]
y = 2.25
for n in notes:
    add_rect(s, 1.0, y + 0.03, 0.14, 0.14, fill=ACCENT)
    add_text(s, 1.35, y - 0.06, 10.8, 0.9, n, size=14, color=hx('#E3EEF4'), ls=1.15)
    y += 1.02

out = "/Users/parndoungjai/Desktop/claude jun 18/May_2026_Engagement_Top5.pptx"
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
